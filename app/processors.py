import base64
from io import BytesIO
from typing import List, Tuple, Dict, Any

from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook

MAX_PDF_PAGES = 20
MAX_TEXT_CHARS = 12000
MAX_IMAGE_BYTES = 4 * 1024 * 1024
MAX_TEXT_FILE_BYTES = 200_000

CODE_EXTS = {
    ".asm", ".s", ".txt", ".md",
    ".css", ".html", ".js", ".ts", ".tsx", ".jsx",
    ".c", ".h", ".cpp", ".hpp", ".cc",
    ".py", ".java", ".kt", ".swift", ".go", ".rs",
    ".rb", ".php", ".sql", ".sh", ".bat", ".ps1",
    ".json", ".yaml", ".yml", ".toml", ".ini",
}


async def build_parts_from_files(files) -> Tuple[List[Dict[str, Any]], Dict[str, Any]]:
    parts: List[Dict[str, Any]] = []
    stats = {"files": 0, "pdf_pages": 0, "text_chars": 0, "images": 0}

    for f in files:
        stats["files"] += 1
        content = await f.read()
        name = (f.filename or "").lower()
        ctype = (f.content_type or "").lower()

        if ctype.startswith("image/") or name.endswith((".png", ".jpg", ".jpeg", ".webp")):
            if len(content) > MAX_IMAGE_BYTES:
                content = content[:MAX_IMAGE_BYTES]
            b64 = base64.b64encode(content).decode()
            parts.append({
                "inline_data": {"mime_type": ctype or "image/png", "data": b64}
            })
            stats["images"] += 1
            continue

        if ctype == "application/pdf" or name.endswith(".pdf"):
            try:
                reader = PdfReader(BytesIO(content))
                text_chunks = []
                for i, page in enumerate(reader.pages[:MAX_PDF_PAGES]):
                    page_text = page.extract_text() or ""
                    text_chunks.append(page_text)
                    stats["pdf_pages"] += 1
                text = "\n".join(text_chunks)
                text = text[:MAX_TEXT_CHARS]
                stats["text_chars"] += len(text)
                if text.strip():
                    parts.append({"text": f"[PDF:{f.filename}]\n{text}"})
            except Exception:
                continue

        # Generic text/code files
        for ext in CODE_EXTS:
            if name.endswith(ext):
                try:
                    text = content[:MAX_TEXT_FILE_BYTES].decode("utf-8", errors="ignore")
                    text = text[:MAX_TEXT_CHARS]
                    stats["text_chars"] += len(text)
                    if text.strip():
                        parts.append({"text": f"[FILE:{f.filename}]\n{text}"})
                except Exception:
                    pass
                break

        if name.endswith(".docx"):
            try:
                doc = Document(BytesIO(content))
                text = "\n".join(p.text for p in doc.paragraphs)
                text = text[:MAX_TEXT_CHARS]
                stats["text_chars"] += len(text)
                if text.strip():
                    parts.append({"text": f"[DOCX:{f.filename}]\n{text}"})
            except Exception:
                continue

        if name.endswith((".pptx", ".ppt")):
            try:
                prs = Presentation(BytesIO(content))
                text_runs = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text_runs.append(shape.text)
                text = "\n".join(text_runs)
                text = text[:MAX_TEXT_CHARS]
                stats["text_chars"] += len(text)
                if text.strip():
                    parts.append({"text": f"[PPTX:{f.filename}]\n{text}"})
            except Exception:
                continue

        if name.endswith((".xlsx", ".xlsm", ".xltx", ".xltm")):
            try:
                wb = load_workbook(BytesIO(content), data_only=True)
                rows = []
                for ws in wb.worksheets:
                    rows.append(f"[Sheet:{ws.title}]")
                    for row in ws.iter_rows(values_only=True):
                        row_text = " | ".join("" if v is None else str(v) for v in row)
                        if row_text.strip():
                            rows.append(row_text)
                        if sum(len(r) for r in rows) > MAX_TEXT_CHARS:
                            break
                    if sum(len(r) for r in rows) > MAX_TEXT_CHARS:
                        break
                text = "\n".join(rows)[:MAX_TEXT_CHARS]
                stats["text_chars"] += len(text)
                if text.strip():
                    parts.append({"text": f"[XLSX:{f.filename}]\n{text}"})
            except Exception:
                continue

    return parts, stats
